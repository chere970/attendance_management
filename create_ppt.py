from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create presentation
prs = Presentation()

# Slide 1: Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Attendance Management System"
subtitle.text = "Internship Project Presentation\nDeveloped by [Your Name]\nDate: [Current Date]"

# Slide 2: Introduction
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Introduction'
tf = body_shape.text_frame
tf.text = 'Project Overview:'

p = tf.add_paragraph()
p.text = '• A comprehensive web-based attendance management system'
p.level = 1

p = tf.add_paragraph()
p.text = '• Designed for organizations to track employee attendance efficiently'
p.level = 1

p = tf.add_paragraph()
p.text = '• Features biometric authentication and request management'
p.level = 1

p = tf.add_paragraph()
p.text = '• Built as a full-stack application during internship'
p.level = 1

# Slide 3: Project Objectives
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Project Objectives'
tf = body_shape.text_frame
tf.text = 'The main objectives of this project were:'

p = tf.add_paragraph()
p.text = '• Develop a secure and user-friendly attendance system'
p.level = 1

p = tf.add_paragraph()
p.text = '• Implement biometric authentication for check-in/check-out'
p.level = 1

p = tf.add_paragraph()
p.text = '• Create admin dashboard for employee and attendance management'
p.level = 1

p = tf.add_paragraph()
p.text = '• Build a request system for leave and other absences'
p.level = 1

p = tf.add_paragraph()
p.text = '• Ensure responsive design and modern UI/UX'
p.level = 1

# Slide 4: Tech Stack
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Technology Stack'
tf = body_shape.text_frame
tf.text = 'Frontend:'

p = tf.add_paragraph()
p.text = '• Next.js 14 (React Framework)'
p.level = 1

p = tf.add_paragraph()
p.text = '• TypeScript for type safety'
p.level = 1

p = tf.add_paragraph()
p.text = '• Tailwind CSS for styling'
p.level = 1

p = tf.add_paragraph()
p.text = '• Shadcn/ui component library'
p.level = 1

p = tf.add_paragraph()
p.text = 'Backend:'
p.level = 0

p = tf.add_paragraph()
p.text = '• Node.js with Express.js'
p.level = 1

p = tf.add_paragraph()
p.text = '• Prisma ORM with MongoDB'
p.level = 1

p = tf.add_paragraph()
p.text = '• JWT for authentication'
p.level = 1

p = tf.add_paragraph()
p.text = '• Multer for file uploads'
p.level = 1

# Slide 5: System Architecture
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'System Architecture'
tf = body_shape.text_frame
tf.text = 'The system follows a client-server architecture:'

p = tf.add_paragraph()
p.text = '• Frontend: Single Page Application (SPA) built with Next.js'
p.level = 1

p = tf.add_paragraph()
p.text = '• Backend: RESTful API server with Express.js'
p.level = 1

p = tf.add_paragraph()
p.text = '• Database: MongoDB with Prisma ORM'
p.level = 1

p = tf.add_paragraph()
p.text = '• Authentication: JWT tokens stored in localStorage'
p.level = 1

p = tf.add_paragraph()
p.text = '• File Storage: Local server storage for employee photos'
p.level = 1

# Slide 6: Database Schema
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Database Schema'
tf = body_shape.text_frame
tf.text = 'Key Models:'

p = tf.add_paragraph()
p.text = '• Employee: User details, credentials, department, photo, fingerprint'
p.level = 1

p = tf.add_paragraph()
p.text = '• Attendance: Daily check-in/check-out records'
p.level = 1

p = tf.add_paragraph()
p.text = '• Request: Leave requests, approvals, and comments'
p.level = 1

p = tf.add_paragraph()
p.text = 'Relationships:'
p.level = 0

p = tf.add_paragraph()
p.text = '• Employee has many Attendance records'
p.level = 1

p = tf.add_paragraph()
p.text = '• Employee has many Requests'
p.level = 1

p = tf.add_paragraph()
p.text = '• Requests can be approved by Admins'
p.level = 1

# Slide 7: Key Features
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Key Features'
tf = body_shape.text_frame
tf.text = 'Employee Features:'

p = tf.add_paragraph()
p.text = '• Biometric check-in/check-out with animated simulation'
p.level = 1

p = tf.add_paragraph()
p.text = '• Submit leave requests (vacation, sick leave, etc.)'
p.level = 1

p = tf.add_paragraph()
p.text = '• View attendance history and request status'
p.level = 1

p = tf.add_paragraph()
p.text = '• Profile management with photo upload'
p.level = 1

p = tf.add_paragraph()
p.text = 'Admin Features:'
p.level = 0

p = tf.add_paragraph()
p.text = '• Manage employees (add, edit, delete)'
p.level = 1

p = tf.add_paragraph()
p.text = '• Approve/reject employee requests'
p.level = 1

p = tf.add_paragraph()
p.text = '• View attendance summaries and working hours'
p.level = 1

p = tf.add_paragraph()
p.text = '• Dashboard with employee overview'
p.level = 1

# Slide 8: User Interface Screenshots
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'User Interface'
tf = body_shape.text_frame
tf.text = 'Modern and Responsive Design:'

p = tf.add_paragraph()
p.text = '• Clean, professional interface using Tailwind CSS'
p.level = 1

p = tf.add_paragraph()
p.text = '• Responsive design for desktop and mobile devices'
p.level = 1

p = tf.add_paragraph()
p.text = '• Intuitive navigation with role-based access'
p.level = 1

p = tf.add_paragraph()
p.text = '• Interactive components with loading states and animations'
p.level = 1

p = tf.add_paragraph()
p.text = 'Key Screens:'
p.level = 0

p = tf.add_paragraph()
p.text = '• Login/Signup pages'
p.level = 1

p = tf.add_paragraph()
p.text = '• Employee attendance page with biometric scanner'
p.level = 1

p = tf.add_paragraph()
p.text = '• Admin dashboard with employee management'
p.level = 1

p = tf.add_paragraph()
p.text = '• Request submission and history pages'
p.level = 1

# Slide 9: Implementation Challenges
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Implementation Challenges'
tf = body_shape.text_frame
tf.text = 'Challenges Faced and Solutions:'

p = tf.add_paragraph()
p.text = '• Biometric Authentication: Simulated with animated UI since hardware not available'
p.level = 1

p = tf.add_paragraph()
p.text = '• State Management: Used React hooks and local state for attendance status'
p.level = 1

p = tf.add_paragraph()
p.text = '• File Uploads: Implemented secure photo upload with Multer middleware'
p.level = 1

p = tf.add_paragraph()
p.text = '• Role-based Access: JWT tokens with role verification on both client and server'
p.level = 1

p = tf.add_paragraph()
p.text = '• Real-time Updates: Polling for attendance status updates'
p.level = 1

# Slide 10: Learning Outcomes
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Learning Outcomes'
tf = body_shape.text_frame
tf.text = 'Skills Developed During Internship:'

p = tf.add_paragraph()
p.text = '• Full-stack development with modern JavaScript frameworks'
p.level = 1

p = tf.add_paragraph()
p.text = '• Database design and ORM usage (Prisma with MongoDB)'
p.level = 1

p = tf.add_paragraph()
p.text = '• RESTful API design and implementation'
p.level = 1

p = tf.add_paragraph()
p.text = '• Authentication and authorization patterns'
p.level = 1

p = tf.add_paragraph()
p.text = '• UI/UX design principles and responsive development'
p.level = 1

p = tf.add_paragraph()
p.text = '• Version control and collaborative development'
p.level = 1

# Slide 11: Future Enhancements
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Future Enhancements'
tf = body_shape.text_frame
tf.text = 'Potential Improvements:'

p = tf.add_paragraph()
p.text = '• Real biometric hardware integration'
p.level = 1

p = tf.add_paragraph()
p.text = '• Real-time notifications and alerts'
p.level = 1

p = tf.add_paragraph()
p.text = '• Advanced reporting and analytics dashboard'
p.level = 1

p = tf.add_paragraph()
p.text = '• Mobile application development'
p.level = 1

p = tf.add_paragraph()
p.text = '• Integration with HR systems and payroll'
p.level = 1

p = tf.add_paragraph()
p.text = '• Multi-language support and accessibility features'
p.level = 1

# Slide 12: Conclusion
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Conclusion'
tf = body_shape.text_frame
tf.text = 'Project Summary:'

p = tf.add_paragraph()
p.text = '• Successfully developed a comprehensive attendance management system'
p.level = 1

p = tf.add_paragraph()
p.text = '• Demonstrated proficiency in full-stack web development'
p.level = 1

p = tf.add_paragraph()
p.text = '• Applied modern development practices and technologies'
p.level = 1

p = tf.add_paragraph()
p.text = '• Created a scalable and maintainable codebase'
p.level = 1

p = tf.add_paragraph()
p.text = '• Gained valuable experience in real-world application development'
p.level = 1

# Slide 13: Q&A
slide = prs.slides.add_slide(prs.slide_layouts[1])
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = 'Q&A'
tf = body_shape.text_frame
tf.text = 'Thank you for your attention!'

p = tf.add_paragraph()
p.text = 'Questions and Discussion'
p.level = 1

p = tf.add_paragraph()
p.text = 'Contact Information:'
p.level = 1

p = tf.add_paragraph()
p.text = '• Email: [your.email@example.com]'
p.level = 2

p = tf.add_paragraph()
p.text = '• GitHub: [github.com/yourusername]'
p.level = 2

p = tf.add_paragraph()
p.text = '• LinkedIn: [linkedin.com/in/yourprofile]'
p.level = 2

# Save the presentation
prs.save('attendance_management_presentation.pptx')
print("Presentation created successfully: attendance_management_presentation.pptx")
